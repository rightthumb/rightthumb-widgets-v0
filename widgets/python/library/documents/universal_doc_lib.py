#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
UniversalDocLib
===============

A document-focused Python library that normalizes many document types into:

1. Markdown
2. Block trees with heading paths
3. Tables / rows where possible
4. Metadata

This file is intentionally heavy on comments because the goal is not just to
work today, but to be easy to build out later.

PRIMARY DESIGN GOAL
-------------------
Everything routes to the same downstream shape no matter the source type.

That means:
- HTML headings
- DOCX styled headings
- PPTX slide titles
- XLSX sheet names / tabular ranges
- PDF inferred headings

...all become a common internal structure.

CURRENT PRIMARY OUTPUT
----------------------
- to_markdown({...})
- to_blocks({...})
- to_tables({...})
- to_rows({...})
- info({...})

SUPPORTED / TARGETED TYPES
--------------------------
Documents:
	pdf
	doc
	docx
	rtf
	odt
	txt
	html / htm / xhtml
	gdoc
	epub
	md / markdown
	tex / latex
	xps

Spreadsheets:
	xls
	xlsx
	csv
	ods
	tsv
	gsheet

Presentations:
	ppt
	pptx
	odp
	gslides

REALISTIC SUPPORT MODEL
-----------------------
Direct adapters:
	- pdf      -> pdfminer.six (+ optional pdfplumber for tables)
	- docx     -> python-docx
	- rtf      -> striprtf
	- html     -> beautifulsoup4
	- xlsx     -> openpyxl
	- pptx     -> python-pptx
	- epub     -> EbookLib
	- latex    -> pylatexenc
	- txt/md   -> native
	- odt/ods/odp -> zip+xml baseline parser

Fallback / future conversion bucket:
	- doc
	- xls
	- ppt
	- xps
	- gdoc
	- gsheet
	- gslides

WHY BLOCKS MATTER
-----------------
Later you want:
- nth heading paths
- vectors by heading path
- vectors by table header signature
- learned parsing profiles
- invoice/table subtotal handling
- HTML regeneration
- advanced markdown regeneration
- section-aware search

So every adapter tries to emit blocks like:

{
	"type": "heading" | "paragraph" | "list" | "list_item" | "table" | "link" | "page_break",
	"text": "Some text",
	"level": 2,
	"path": ["Top Heading", "Sub Heading"],
	"attrs": {...}
}

INSTALL
-------
Minimum useful install:
	pip install beautifulsoup4

Recommended:
	pip install pdfminer.six pdfplumber python-docx openpyxl python-pptx
	pip install EbookLib pylatexenc striprtf rich

Optional:
	pip install pymupdf4llm

Notes:
- This file does NOT require every dependency to exist.
- It gracefully degrades and leaves hooks for future fallbacks.

FUTURE TODO IDEAS
-----------------
- better PDF heading inference using font size / weight / spacing / position
- learned profile store keyed by document type + table signature + title pattern
- converter hooks for doc/xls/ppt/xps/google workspace
- better ODT/ODS/ODP namespace parsing
- footnotes / comments / notes extraction
- invoice summary-line detection outside formal tables
- richer DOCX list detection
- true HTML->Markdown converter option
- YAML / JSON export of block tree
"""

import os
import re
import io
import csv
import json
import zipfile
import shutil
import tempfile
import subprocess
import xml.etree.ElementTree as ET
from collections import defaultdict
from typing import Any, Dict, List, Optional, Tuple


class UniversalDocLib:
	"""
	Main document abstraction class.

	All public methods take a single dict.
	The constructor also takes a single dict.

	Example:
		doc = UniversalDocLib({"path": "file.docx"})
		md = doc.to_markdown({})
		blocks = doc.to_blocks({})
		tables = doc.to_tables({})
	"""

	def __init__(self, cfg: Optional[Dict[str, Any]] = None):
		self.cfg = self._defaults(cfg, {
			"path": None,
			"kind": "auto",
			"encoding": "utf-8",
			"errors": "replace",
			"debug": False,
			"cache": True,
			"fallback_convert": False,
			"tmp_dir": None,
			"pdf_table_mode": "prominent",   # prominent | all | none
			"markdown_tables": True,
		})

		self.state = {
			"path": self.cfg.get("path"),
			"kind": None,
			"meta": {},
			"text": None,
			"markdown": None,
			"blocks": None,
			"tables": None,
			"rows": None,
		}

		self.state["kind"] = self._infer_kind({
			"path": self.cfg.get("path"),
			"kind": self.cfg.get("kind", "auto"),
		})

	# =========================================================================
	# Public API
	# =========================================================================

	def info(self, cfg: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
		"""
		Return metadata / state summary.

		Good for debugging what the adapter thought the file type was.
		"""
		merged = self._merge_cfg(cfg)
		path = merged.get("path") or self.state.get("path")
		kind = self._effective_kind(merged)

		return {
			"path": path,
			"kind": kind,
			"has_text": self.state.get("text") is not None,
			"has_markdown": self.state.get("markdown") is not None,
			"has_blocks": self.state.get("blocks") is not None,
			"has_tables": self.state.get("tables") is not None,
			"meta": self.state.get("meta", {}),
		}

	def to_text(self, cfg: Optional[Dict[str, Any]] = None) -> str:
		"""
		Plain text extraction.

		This is not the primary target of the system anymore, but it is still
		useful as a fallback and for debugging.
		"""
		merged = self._merge_cfg(cfg, {"cache": True})

		if merged["cache"] and self.state.get("text") is not None:
			return self.state["text"]

		kind = self._effective_kind(merged)

		if kind in ("txt", "text", "md", "markdown", "json", "xml"):
			text = self._read_text_file(merged)

		elif kind in ("html", "htm", "xhtml"):
			text = self._html_to_text(merged)

		elif kind == "pdf":
			text = self._pdf_to_text(merged)

		elif kind == "docx":
			text = self._docx_to_text(merged)

		elif kind == "rtf":
			text = self._rtf_to_text(merged)

		elif kind == "odt":
			text = self._odt_to_text(merged)

		elif kind in ("xlsx", "xlsm", "xltx", "xltm"):
			text = self._xlsx_to_text(merged)

		elif kind == "ods":
			text = self._ods_to_text(merged)

		elif kind == "csv":
			text = self._csv_to_text({**merged, "delimiter": ","})

		elif kind == "tsv":
			text = self._csv_to_text({**merged, "delimiter": "\t"})

		elif kind == "pptx":
			text = self._pptx_to_text(merged)

		elif kind == "odp":
			text = self._odp_to_text(merged)

		elif kind == "epub":
			text = self._epub_to_text(merged)

		elif kind in ("tex", "latex"):
			text = self._latex_to_text(merged)

		else:
			text = self._fallback_to_text(merged)

		if merged["cache"]:
			self.state["text"] = text

		return text

	def to_markdown(self, cfg: Optional[Dict[str, Any]] = None) -> str:
		"""
		Primary user-facing output.

		Convert document into normalized Markdown with:
		- headings
		- paragraphs
		- lists
		- tables
		- page or slide boundaries where useful

		Most adapters build blocks first, then render blocks to markdown.
		"""
		merged = self._merge_cfg(cfg, {"cache": True})

		if merged["cache"] and self.state.get("markdown") is not None:
			return self.state["markdown"]

		blocks = self.to_blocks(merged)
		md = self._blocks_to_markdown({"blocks": blocks})

		if merged["cache"]:
			self.state["markdown"] = md

		return md

	def to_blocks(self, cfg: Optional[Dict[str, Any]] = None) -> List[Dict[str, Any]]:
		"""
		Core normalized output.

		Every adapter should try to return block dictionaries with heading paths.

		This is the most important output for future vectoring / profiling / HTML.
		"""
		merged = self._merge_cfg(cfg, {"cache": True})

		if merged["cache"] and self.state.get("blocks") is not None:
			return self.state["blocks"]

		kind = self._effective_kind(merged)

		if kind in ("txt", "text"):
			blocks = self._text_to_blocks({"text": self._read_text_file(merged), "kind": kind})

		elif kind in ("md", "markdown"):
			blocks = self._markdown_text_to_blocks({"text": self._read_text_file(merged), "kind": kind})

		elif kind in ("html", "htm", "xhtml"):
			blocks = self._html_to_blocks(merged)

		elif kind == "pdf":
			blocks = self._pdf_to_blocks(merged)

		elif kind == "docx":
			blocks = self._docx_to_blocks(merged)

		elif kind == "rtf":
			text = self._rtf_to_text(merged)
			blocks = self._text_to_blocks({"text": text, "kind": kind})

		elif kind == "odt":
			blocks = self._odt_to_blocks(merged)

		elif kind in ("xlsx", "xlsm", "xltx", "xltm"):
			blocks = self._xlsx_to_blocks(merged)

		elif kind == "ods":
			blocks = self._ods_to_blocks(merged)

		elif kind in ("csv", "tsv"):
			blocks = self._delimited_to_blocks(merged)

		elif kind == "pptx":
			blocks = self._pptx_to_blocks(merged)

		elif kind == "odp":
			blocks = self._odp_to_blocks(merged)

		elif kind == "epub":
			blocks = self._epub_to_blocks(merged)

		elif kind in ("tex", "latex"):
			text = self._latex_to_text(merged)
			blocks = self._text_to_blocks({"text": text, "kind": kind})

		else:
			text = self._fallback_to_text(merged)
			blocks = self._text_to_blocks({"text": text, "kind": kind})

		if merged["cache"]:
			self.state["blocks"] = blocks

		return blocks

	def to_tables(self, cfg: Optional[Dict[str, Any]] = None) -> List[Dict[str, Any]]:
		"""
		Return normalized table objects.

		Table object shape:
		{
			"path": [...],
			"headers": [...],
			"rows": [[...], [...]],
			"attrs": {...}
		}
		"""
		merged = self._merge_cfg(cfg, {"cache": True})

		if merged["cache"] and self.state.get("tables") is not None:
			return self.state["tables"]

		blocks = self.to_blocks(merged)
		tables = []

		for block in blocks:
			if block.get("type") == "table":
				tables.append({
					"path": block.get("path", []),
					"headers": block.get("headers", []),
					"rows": block.get("rows", []),
					"attrs": block.get("attrs", {}),
				})

		if merged["cache"]:
			self.state["tables"] = tables

		return tables

	def to_rows(self, cfg: Optional[Dict[str, Any]] = None) -> List[Dict[str, Any]]:
		"""
		Flatten tables into list-of-dict rows.

		This is mostly for convenient downstream processing.
		"""
		merged = self._merge_cfg(cfg, {"cache": True})

		if merged["cache"] and self.state.get("rows") is not None:
			return self.state["rows"]

		tables = self.to_tables(merged)
		rows_out = []

		for ti, table in enumerate(tables):
			headers = table.get("headers", [])
			for row in table.get("rows", []):
				obj = {}
				for i, header in enumerate(headers):
					key = self._clean_header(header) or f"col_{i+1}"
					obj[key] = row[i] if i < len(row) else ""
				obj["_table_index"] = ti
				obj["_path"] = table.get("path", [])
				rows_out.append(obj)

		if merged["cache"]:
			self.state["rows"] = rows_out

		return rows_out

	# =========================================================================
	# HTML adapter
	# =========================================================================

	def _html_to_blocks(self, cfg: Dict[str, Any]) -> List[Dict[str, Any]]:
		"""
		HTML is the cleanest structured source in this whole stack.

		We can map:
		- h1..h6 -> heading blocks
		- p -> paragraph
		- ul/ol/li -> list/list_item
		- table -> table
		- a -> link
		"""
		html = self._read_text_file(cfg)
		if not html:
			return []

		try:
			from bs4 import BeautifulSoup, NavigableString, Tag
		except Exception:
			# Fallback: just strip tags and make paragraphs.
			text = re.sub(r"<[^>]+>", " ", html)
			return self._text_to_blocks({"text": text, "kind": "html"})

		soup = BeautifulSoup(html, "html.parser")
		body = soup.body if soup.body else soup

		blocks = []
		heading_stack = []

		def set_heading(level: int, text: str) -> List[str]:
			while len(heading_stack) >= level:
				heading_stack.pop()
			heading_stack.append(text)
			return list(heading_stack)

		def current_path() -> List[str]:
			return list(heading_stack)

		def handle_table(table_tag):
			matrix = []
			for tr in table_tag.find_all("tr"):
				cells = tr.find_all(["th", "td"])
				matrix.append([c.get_text(" ", strip=True) for c in cells])

			headers, rows = self._matrix_split_header_rows(matrix)
			blocks.append({
				"type": "table",
				"path": current_path(),
				"headers": headers,
				"rows": rows,
				"attrs": {"source_kind": "html"},
			})

		def walk(node):
			if isinstance(node, NavigableString):
				return

			if not isinstance(node, Tag):
				return

			name = node.name.lower()

			if re.fullmatch(r"h[1-6]", name):
				level = int(name[1])
				text = node.get_text(" ", strip=True)
				if text:
					path = set_heading(level, text)
					blocks.append({
						"type": "heading",
						"text": text,
						"level": level,
						"path": path,
						"attrs": {"source_kind": "html"},
					})
				return

			if name == "p":
				text = node.get_text(" ", strip=True)
				if text:
					blocks.append({
						"type": "paragraph",
						"text": text,
						"path": current_path(),
						"attrs": {"source_kind": "html"},
					})
				return

			if name in ("ul", "ol"):
				items = node.find_all("li", recursive=False)
				if items:
					blocks.append({
						"type": "list",
						"ordered": name == "ol",
						"path": current_path(),
						"attrs": {"source_kind": "html"},
					})
					for li in items:
						text = li.get_text(" ", strip=True)
						if text:
							blocks.append({
								"type": "list_item",
								"text": text,
								"path": current_path(),
								"attrs": {"source_kind": "html"},
							})
				return

			if name == "table":
				handle_table(node)
				return

			if name == "a":
				text = node.get_text(" ", strip=True)
				href = node.get("href")
				if text or href:
					blocks.append({
						"type": "link",
						"text": text or href or "",
						"path": current_path(),
						"attrs": {"href": href, "source_kind": "html"},
					})
				return

			for child in node.children:
				walk(child)

		for child in body.children:
			walk(child)

		return blocks

	def _html_to_text(self, cfg: Dict[str, Any]) -> str:
		blocks = self._html_to_blocks(cfg)
		out = []

		for block in blocks:
			if block["type"] in ("heading", "paragraph", "list_item", "link"):
				out.append(block.get("text", ""))
			elif block["type"] == "table":
				out.append(self._table_block_to_plain_text({"block": block}))
			elif block["type"] == "list":
				pass

		return "\n".join(x for x in out if x).strip()

	# =========================================================================
	# DOCX adapter
	# =========================================================================

	def _docx_to_blocks(self, cfg: Dict[str, Any]) -> List[Dict[str, Any]]:
		"""
		DOCX is a great structured source because paragraph styles exist.

		We use:
		- Heading N styles -> heading blocks
		- normal paragraphs -> paragraph blocks
		- tables -> table blocks

		Future:
		- better list detection
		- images, captions, comments, footnotes
		"""
		path = cfg.get("path") or self.state.get("path")
		if not path:
			return []

		try:
			import docx
		except Exception:
			text = self._fallback_to_text(cfg)
			return self._text_to_blocks({"text": text, "kind": "docx"})

		try:
			d = docx.Document(path)
		except Exception:
			text = self._fallback_to_text(cfg)
			return self._text_to_blocks({"text": text, "kind": "docx"})

		blocks = []
		heading_stack = []

		def set_heading(level: int, text: str) -> List[str]:
			while len(heading_stack) >= level:
				heading_stack.pop()
			heading_stack.append(text)
			return list(heading_stack)

		def current_path() -> List[str]:
			return list(heading_stack)

		# Paragraphs
		for p in d.paragraphs:
			text = p.text.strip()
			if not text:
				continue

			style_name = ""
			try:
				style_name = p.style.name or ""
			except Exception:
				style_name = ""

			m = re.match(r"Heading\s+(\d+)", style_name, re.I)
			if m:
				level = int(m.group(1))
				path_now = set_heading(level, text)
				blocks.append({
					"type": "heading",
					"text": text,
					"level": level,
					"path": path_now,
					"attrs": {
						"source_kind": "docx",
						"style_name": style_name,
					},
				})
			else:
				blocks.append({
					"type": "paragraph",
					"text": text,
					"path": current_path(),
					"attrs": {
						"source_kind": "docx",
						"style_name": style_name,
					},
				})

		# Tables
		for table in d.tables:
			matrix = []
			for row in table.rows:
				matrix.append([cell.text.strip() for cell in row.cells])

			headers, rows = self._matrix_split_header_rows(matrix)
			blocks.append({
				"type": "table",
				"path": current_path(),
				"headers": headers,
				"rows": rows,
				"attrs": {"source_kind": "docx"},
			})

		return blocks

	def _docx_to_text(self, cfg: Dict[str, Any]) -> str:
		blocks = self._docx_to_blocks(cfg)
		return self._blocks_to_plain_text({"blocks": blocks})

	# =========================================================================
	# PDF adapter
	# =========================================================================

	def _pdf_to_blocks(self, cfg: Dict[str, Any]) -> List[Dict[str, Any]]:
		"""
		PDF is the hardest one.

		Current practical strategy:
		- get text from pdfminer.six
		- create paragraph-ish blocks
		- optionally extract tables with pdfplumber
		- add very conservative heading inference

		Because PDF stores positioned characters rather than semantic structure,
		any heading detection here is heuristic.
		"""
		text = self._pdf_to_text(cfg)
		blocks = []

		# Conservative text -> paragraph/heading heuristic.
		# This is intentionally simple for now and easy to replace later.
		paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text or "") if p.strip()]
		heading_stack = []

		def set_heading(level: int, text_val: str) -> List[str]:
			while len(heading_stack) >= level:
				heading_stack.pop()
			heading_stack.append(text_val)
			return list(heading_stack)

		def current_path() -> List[str]:
			return list(heading_stack)

		for p in paragraphs:
			lines = [x.strip() for x in p.splitlines() if x.strip()]
			candidate = " ".join(lines)

			if self._looks_like_heading({"text": candidate}):
				level = 1 if len(heading_stack) == 0 else min(len(heading_stack) + 1, 6)
				path = set_heading(level, candidate)
				blocks.append({
					"type": "heading",
					"text": candidate,
					"level": level,
					"path": path,
					"attrs": {"source_kind": "pdf", "heuristic": True},
				})
			else:
				blocks.append({
					"type": "paragraph",
					"text": candidate,
					"path": current_path(),
					"attrs": {"source_kind": "pdf"},
				})

		# Tables
		table_mode = cfg.get("pdf_table_mode", self.cfg.get("pdf_table_mode", "prominent"))
		if table_mode != "none":
			table_blocks = self._pdf_extract_table_blocks(cfg)
			blocks.extend(table_blocks)

		return blocks

	def _pdf_to_text(self, cfg: Dict[str, Any]) -> str:
		path = cfg.get("path") or self.state.get("path")
		if not path:
			return ""

		try:
			from pdfminer.high_level import extract_text
			return extract_text(path) or ""
		except Exception:
			return self._fallback_to_text(cfg)

	def _pdf_extract_table_blocks(self, cfg: Dict[str, Any]) -> List[Dict[str, Any]]:
		"""
		Extract PDF tables using pdfplumber if present.

		Modes:
		- prominent: just strongest candidate
		- all: all candidates
		"""
		path = cfg.get("path") or self.state.get("path")
		if not path:
			return []

		try:
			import pdfplumber
		except Exception:
			return []

		candidates = []
		try:
			with pdfplumber.open(path) as pdf:
				for page_num, page in enumerate(pdf.pages, start=1):
					for tbl in (page.extract_tables() or []):
						if not tbl:
							continue

						row_count = len(tbl)
						col_count = max((len(r) for r in tbl if r), default=0)
						non_empty = sum(1 for r in tbl for c in (r or []) if c not in (None, "", " "))
						score = (row_count * 3) + (col_count * 2) + non_empty

						candidates.append({
							"page": page_num,
							"matrix": tbl,
							"_score": score,
						})
		except Exception:
			return []

		if not candidates:
			return []

		mode = cfg.get("pdf_table_mode", self.cfg.get("pdf_table_mode", "prominent"))
		selected = candidates if mode == "all" else [sorted(candidates, key=lambda x: x["_score"], reverse=True)[0]]

		blocks = []
		for item in selected:
			headers, rows = self._matrix_split_header_rows(item["matrix"])
			blocks.append({
				"type": "table",
				"path": [],
				"headers": headers,
				"rows": rows,
				"attrs": {
					"source_kind": "pdf",
					"page": item["page"],
					"score": item["_score"],
				},
			})

		return blocks

	# =========================================================================
	# XLSX / CSV / TSV / ODS adapters
	# =========================================================================

	def _xlsx_to_blocks(self, cfg: Dict[str, Any]) -> List[Dict[str, Any]]:
		path = cfg.get("path") or self.state.get("path")
		if not path:
			return []

		try:
			from openpyxl import load_workbook
		except Exception:
			text = self._fallback_to_text(cfg)
			return self._text_to_blocks({"text": text, "kind": "xlsx"})

		try:
			wb = load_workbook(path, data_only=True, read_only=True)
		except Exception:
			text = self._fallback_to_text(cfg)
			return self._text_to_blocks({"text": text, "kind": "xlsx"})

		blocks = []

		for ws in wb.worksheets:
			sheet_name = ws.title
			path_now = [sheet_name]

			blocks.append({
				"type": "heading",
				"text": sheet_name,
				"level": 1,
				"path": path_now,
				"attrs": {"source_kind": "xlsx", "sheet": sheet_name},
			})

			matrix = []
			for row in ws.iter_rows(values_only=True):
				vals = ["" if v is None else str(v) for v in row]
				if any(str(v).strip() for v in vals):
					matrix.append(vals)

			if matrix:
				headers, rows = self._matrix_split_header_rows(matrix)
				blocks.append({
					"type": "table",
					"path": path_now,
					"headers": headers,
					"rows": rows,
					"attrs": {"source_kind": "xlsx", "sheet": sheet_name},
				})

		return blocks

	def _xlsx_to_text(self, cfg: Dict[str, Any]) -> str:
		return self._blocks_to_plain_text({"blocks": self._xlsx_to_blocks(cfg)})

	def _delimited_to_blocks(self, cfg: Dict[str, Any]) -> List[Dict[str, Any]]:
		delimiter = "," if self._effective_kind(cfg) == "csv" else "\t"
		path = cfg.get("path") or self.state.get("path")
		if not path:
			return []

		rows = []
		try:
			with open(path, "r", encoding=cfg.get("encoding", "utf-8"), errors=cfg.get("errors", "replace"), newline="") as f:
				reader = csv.reader(f, delimiter=delimiter)
				for row in reader:
					rows.append([str(x) for x in row])
		except Exception:
			return []

		if not rows:
			return []

		headers, body = self._matrix_split_header_rows(rows)
		title = os.path.basename(path)

		return [
			{
				"type": "heading",
				"text": title,
				"level": 1,
				"path": [title],
				"attrs": {"source_kind": self._effective_kind(cfg)},
			},
			{
				"type": "table",
				"path": [title],
				"headers": headers,
				"rows": body,
				"attrs": {"source_kind": self._effective_kind(cfg)},
			}
		]

	def _csv_to_text(self, cfg: Dict[str, Any]) -> str:
		return self._blocks_to_plain_text({"blocks": self._delimited_to_blocks({**cfg, "kind": "csv"})})

	def _ods_to_blocks(self, cfg: Dict[str, Any]) -> List[Dict[str, Any]]:
		"""
		Lightweight ODS parser from content.xml.

		This is baseline support. It is intentionally simple and replaceable.
		"""
		path = cfg.get("path") or self.state.get("path")
		if not path:
			return []

		try:
			with zipfile.ZipFile(path, "r") as zf:
				raw = zf.read("content.xml")
		except Exception:
			text = self._fallback_to_text(cfg)
			return self._text_to_blocks({"text": text, "kind": "ods"})

		try:
			root = ET.fromstring(raw)
		except Exception:
			text = self._fallback_to_text(cfg)
			return self._text_to_blocks({"text": text, "kind": "ods"})

		ns = {
			"table": "urn:oasis:names:tc:opendocument:xmlns:table:1.0",
			"text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
		}

		blocks = []

		for table_el in root.findall(".//table:table", ns):
			sheet_name = table_el.attrib.get("{urn:oasis:names:tc:opendocument:xmlns:table:1.0}name", "Sheet")
			path_now = [sheet_name]

			blocks.append({
				"type": "heading",
				"text": sheet_name,
				"level": 1,
				"path": path_now,
				"attrs": {"source_kind": "ods", "sheet": sheet_name},
			})

			matrix = []
			for row_el in table_el.findall("./table:table-row", ns):
				vals = []
				for cell_el in row_el.findall("./table:table-cell", ns):
					txts = [t.text or "" for t in cell_el.findall(".//text:p", ns)]
					vals.append(" ".join(x.strip() for x in txts if x is not None).strip())
				if any(v.strip() for v in vals):
					matrix.append(vals)

			if matrix:
				headers, rows = self._matrix_split_header_rows(matrix)
				blocks.append({
					"type": "table",
					"path": path_now,
					"headers": headers,
					"rows": rows,
					"attrs": {"source_kind": "ods", "sheet": sheet_name},
				})

		return blocks

	def _ods_to_text(self, cfg: Dict[str, Any]) -> str:
		return self._blocks_to_plain_text({"blocks": self._ods_to_blocks(cfg)})

	# =========================================================================
	# PPTX / ODP adapters
	# =========================================================================

	def _pptx_to_blocks(self, cfg: Dict[str, Any]) -> List[Dict[str, Any]]:
		"""
		Slide title => heading.
		Other text shapes => paragraphs.
		Tables => table blocks.

		Future:
		- bullet nesting
		- notes
		- placeholder type detection
		"""
		path = cfg.get("path") or self.state.get("path")
		if not path:
			return []

		try:
			from pptx import Presentation
		except Exception:
			text = self._fallback_to_text(cfg)
			return self._text_to_blocks({"text": text, "kind": "pptx"})

		try:
			prs = Presentation(path)
		except Exception:
			text = self._fallback_to_text(cfg)
			return self._text_to_blocks({"text": text, "kind": "pptx"})

		blocks = []

		for idx, slide in enumerate(prs.slides, start=1):
			title = f"Slide {idx}"
			slide_title = None

			# Best-effort title extraction
			try:
				if slide.shapes.title and slide.shapes.title.text.strip():
					slide_title = slide.shapes.title.text.strip()
			except Exception:
				slide_title = None

			if slide_title:
				title = slide_title

			path_now = [title]

			blocks.append({
				"type": "heading",
				"text": title,
				"level": 1,
				"path": path_now,
				"attrs": {"source_kind": "pptx", "slide": idx},
			})

			for shape in slide.shapes:
				# Tables
				if getattr(shape, "has_table", False):
					matrix = []
					for row in shape.table.rows:
						matrix.append([cell.text.strip() for cell in row.cells])

					headers, rows = self._matrix_split_header_rows(matrix)
					blocks.append({
						"type": "table",
						"path": path_now,
						"headers": headers,
						"rows": rows,
						"attrs": {"source_kind": "pptx", "slide": idx},
					})
					continue

				# Text
				if getattr(shape, "has_text_frame", False):
					text = shape.text.strip()
					if text and text != title:
						blocks.append({
							"type": "paragraph",
							"text": text,
							"path": path_now,
							"attrs": {"source_kind": "pptx", "slide": idx},
						})

		return blocks

	def _pptx_to_text(self, cfg: Dict[str, Any]) -> str:
		return self._blocks_to_plain_text({"blocks": self._pptx_to_blocks(cfg)})

	def _odp_to_blocks(self, cfg: Dict[str, Any]) -> List[Dict[str, Any]]:
		"""
		Lightweight fallback based on content.xml text extraction.
		Better ODP semantic parsing can be added later.
		"""
		text = self._odp_to_text(cfg)
		blocks = self._text_to_blocks({"text": text, "kind": "odp"})
		if blocks:
			title = os.path.basename(cfg.get("path") or self.state.get("path") or "Presentation")
			blocks.insert(0, {
				"type": "heading",
				"text": title,
				"level": 1,
				"path": [title],
				"attrs": {"source_kind": "odp"},
			})
		return blocks

	def _odp_to_text(self, cfg: Dict[str, Any]) -> str:
		return self._odf_text_fallback(cfg)

	# =========================================================================
	# ODT / RTF / EPUB / LaTeX adapters
	# =========================================================================

	def _odt_to_blocks(self, cfg: Dict[str, Any]) -> List[Dict[str, Any]]:
		text = self._odt_to_text(cfg)
		return self._text_to_blocks({"text": text, "kind": "odt"})

	def _odt_to_text(self, cfg: Dict[str, Any]) -> str:
		return self._odf_text_fallback(cfg)

	def _rtf_to_text(self, cfg: Dict[str, Any]) -> str:
		path = cfg.get("path") or self.state.get("path")
		if not path:
			return ""

		try:
			from striprtf.striprtf import rtf_to_text
			raw = self._read_text_file(cfg)
			return rtf_to_text(raw)
		except Exception:
			raw = self._read_text_file(cfg)
			raw = re.sub(r"\\[a-zA-Z]+\d* ?", "", raw)
			raw = re.sub(r"[{}]", "", raw)
			return raw.strip()

	def _epub_to_blocks(self, cfg: Dict[str, Any]) -> List[Dict[str, Any]]:
		path = cfg.get("path") or self.state.get("path")
		if not path:
			return []

		try:
			from ebooklib import epub, ITEM_DOCUMENT
			from bs4 import BeautifulSoup
		except Exception:
			text = self._fallback_to_text(cfg)
			return self._text_to_blocks({"text": text, "kind": "epub"})

		try:
			book = epub.read_epub(path)
		except Exception:
			text = self._fallback_to_text(cfg)
			return self._text_to_blocks({"text": text, "kind": "epub"})

		blocks = []

		for idx, item in enumerate(book.get_items(), start=1):
			if item.get_type() != ITEM_DOCUMENT:
				continue

			soup = BeautifulSoup(item.get_content(), "html.parser")
			title = None
			h = soup.find(["h1", "h2"])
			if h:
				title = h.get_text(" ", strip=True)

			if title:
				path_now = [title]
				blocks.append({
					"type": "heading",
					"text": title,
					"level": 1,
					"path": path_now,
					"attrs": {"source_kind": "epub", "chapter_index": idx},
				})
			else:
				path_now = []

			for p in soup.find_all("p"):
				text = p.get_text(" ", strip=True)
				if text:
					blocks.append({
						"type": "paragraph",
						"text": text,
						"path": path_now,
						"attrs": {"source_kind": "epub", "chapter_index": idx},
					})

		return blocks

	def _epub_to_text(self, cfg: Dict[str, Any]) -> str:
		return self._blocks_to_plain_text({"blocks": self._epub_to_blocks(cfg)})

	def _latex_to_text(self, cfg: Dict[str, Any]) -> str:
		raw = self._read_text_file(cfg)
		if not raw:
			return ""

		try:
			from pylatexenc.latex2text import LatexNodes2Text
			return LatexNodes2Text().latex_to_text(raw)
		except Exception:
			text = re.sub(r"\\[a-zA-Z]+\*?(\[[^\]]*\])?(\{[^}]*\})?", " ", raw)
			text = re.sub(r"[{}]", " ", text)
			return re.sub(r"\s+", " ", text).strip()

	# =========================================================================
	# Generic block builders / renderers
	# =========================================================================

	def _text_to_blocks(self, cfg: Dict[str, Any]) -> List[Dict[str, Any]]:
		"""
		Plain text fallback.

		Splits on blank lines into paragraph-ish blocks.
		"""
		text = cfg.get("text", "") or ""
		paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
		blocks = []

		for p in paragraphs:
			blocks.append({
				"type": "paragraph",
				"text": " ".join(line.strip() for line in p.splitlines() if line.strip()),
				"path": [],
				"attrs": {"source_kind": cfg.get("kind", "text")},
			})

		return blocks

	def _markdown_text_to_blocks(self, cfg: Dict[str, Any]) -> List[Dict[str, Any]]:
		"""
		Very lightweight Markdown parser.

		This is only a baseline and intentionally easy to replace.
		"""
		text = cfg.get("text", "") or ""
		lines = text.splitlines()
		blocks = []
		heading_stack = []

		def set_heading(level: int, txt: str) -> List[str]:
			while len(heading_stack) >= level:
				heading_stack.pop()
			heading_stack.append(txt)
			return list(heading_stack)

		for line in lines:
			s = line.strip()
			if not s:
				continue

			m = re.match(r"^(#{1,6})\s+(.*)$", s)
			if m:
				level = len(m.group(1))
				txt = m.group(2).strip()
				path_now = set_heading(level, txt)
				blocks.append({
					"type": "heading",
					"text": txt,
					"level": level,
					"path": path_now,
					"attrs": {"source_kind": "markdown"},
				})
				continue

			if re.match(r"^[-*]\s+", s):
				blocks.append({
					"type": "list_item",
					"text": re.sub(r"^[-*]\s+", "", s),
					"path": list(heading_stack),
					"attrs": {"source_kind": "markdown"},
				})
				continue

			blocks.append({
				"type": "paragraph",
				"text": s,
				"path": list(heading_stack),
				"attrs": {"source_kind": "markdown"},
			})

		return blocks

	def _blocks_to_markdown(self, cfg: Dict[str, Any]) -> str:
		"""
		Normalized renderer from blocks -> Markdown.

		This is where all adapters unify.
		"""
		blocks = cfg.get("blocks", []) or []
		out = []

		for block in blocks:
			t = block.get("type")

			if t == "heading":
				level = max(1, min(int(block.get("level", 1)), 6))
				out.append(("#" * level) + " " + block.get("text", ""))
				out.append("")

			elif t == "paragraph":
				out.append(block.get("text", ""))
				out.append("")

			elif t == "list":
				# container block; actual items printed by list_item blocks
				pass

			elif t == "list_item":
				out.append("- " + block.get("text", ""))
				# keep list blocks tight

			elif t == "link":
				text = block.get("text", "") or block.get("attrs", {}).get("href", "")
				href = block.get("attrs", {}).get("href")
				if href:
					out.append(f"[{text}]({href})")
				else:
					out.append(text)
				out.append("")

			elif t == "table":
				out.append(self._table_block_to_markdown({"block": block}))
				out.append("")

		# normalize blank lines
		rendered = "\n".join(out)
		rendered = re.sub(r"\n{3,}", "\n\n", rendered).strip()
		return rendered

	def _blocks_to_plain_text(self, cfg: Dict[str, Any]) -> str:
		blocks = cfg.get("blocks", []) or []
		out = []

		for block in blocks:
			t = block.get("type")
			if t in ("heading", "paragraph", "list_item", "link"):
				out.append(block.get("text", ""))
			elif t == "table":
				out.append(self._table_block_to_plain_text({"block": block}))

		return "\n".join(x for x in out if x).strip()

	def _table_block_to_markdown(self, cfg: Dict[str, Any]) -> str:
		block = cfg.get("block", {}) or {}
		headers = [str(x) for x in block.get("headers", [])]
		rows = [[str(c) for c in row] for row in block.get("rows", [])]

		if not headers and not rows:
			return ""

		if not headers and rows:
			width = max((len(r) for r in rows), default=0)
			headers = [f"col_{i+1}" for i in range(width)]

		matrix = [headers] + rows
		return self._matrix_to_markdown({"matrix": matrix})

	def _table_block_to_plain_text(self, cfg: Dict[str, Any]) -> str:
		block = cfg.get("block", {}) or {}
		headers = [str(x) for x in block.get("headers", [])]
		rows = [[str(c) for c in row] for row in block.get("rows", [])]
		lines = []

		if headers:
			lines.append(" | ".join(headers))
		for row in rows:
			lines.append(" | ".join(row))

		return "\n".join(lines)

	def _matrix_to_markdown(self, cfg: Dict[str, Any]) -> str:
		matrix = cfg.get("matrix", []) or []
		if not matrix:
			return ""

		width = max((len(r) for r in matrix), default=0)
		norm = []

		for row in matrix:
			r = list(row)
			if len(r) < width:
				r.extend([""] * (width - len(r)))
			norm.append([str(x) if x is not None else "" for x in r])

		header = norm[0]
		body = norm[1:]

		out = []
		out.append("| " + " | ".join(header) + " |")
		out.append("| " + " | ".join(["---"] * len(header)) + " |")
		for row in body:
			out.append("| " + " | ".join(row) + " |")

		return "\n".join(out)

	def _matrix_split_header_rows(self, matrix: List[List[Any]]) -> Tuple[List[str], List[List[str]]]:
		"""
		Turn raw matrix into header row + body rows.

		Current assumption:
		- first row is the header
		- future version can learn header signatures and detect multi-row headers
		"""
		if not matrix:
			return [], []

		width = max((len(r) for r in matrix), default=0)
		norm = []
		for row in matrix:
			r = list(row)
			if len(r) < width:
				r.extend([""] * (width - len(r)))
			norm.append([("" if x is None else str(x).strip()) for x in r])

		header = norm[0]
		rows = norm[1:]

		header = [h if h else f"col_{i+1}" for i, h in enumerate(header)]
		return header, rows

	# =========================================================================
	# Helpers / fallbacks
	# =========================================================================

	def _looks_like_heading(self, cfg: Dict[str, Any]) -> bool:
		"""
		Very conservative heading heuristic for weakly structured sources like PDF.

		Future:
		- font-size-based inference
		- center alignment, numbering, spacing, title detection
		"""
		text = (cfg.get("text") or "").strip()
		if not text:
			return False

		if len(text) > 120:
			return False

		if text.endswith("."):
			return False

		# More title-like than paragraph-like
		if text.isupper() and len(text.split()) <= 10:
			return True

		if len(text.split()) <= 8 and re.search(r"[A-Za-z]", text):
			# Avoid obvious sentence-ish fragments
			if "," not in text and ";" not in text:
				return True

		return False

	def _odf_text_fallback(self, cfg: Dict[str, Any]) -> str:
		path = cfg.get("path") or self.state.get("path")
		if not path:
			return ""

		try:
			with zipfile.ZipFile(path, "r") as zf:
				raw = zf.read("content.xml").decode("utf-8", errors="replace")
			text = re.sub(r"</(text:p|text:h|table:table-row)>", "\n", raw)
			text = re.sub(r"<[^>]+>", " ", text)
			text = re.sub(r"[ \t]+", " ", text)
			text = re.sub(r"\n{3,}", "\n\n", text)
			return text.strip()
		except Exception:
			return self._fallback_to_text(cfg)

	def _read_text_file(self, cfg: Dict[str, Any]) -> str:
		path = cfg.get("path") or self.state.get("path")
		if not path:
			return ""

		try:
			with open(path, "r", encoding=cfg.get("encoding", "utf-8"), errors=cfg.get("errors", "replace")) as f:
				return f.read()
		except Exception:
			return ""

	def _fallback_to_text(self, cfg: Dict[str, Any]) -> str:
		"""
		Optional converter bucket.

		Disabled by default because you asked for the library app first, not a
		bunch of converter assumptions. You can flip fallback_convert=True later.

		Future:
		- LibreOffice headless
		- pandoc
		- google export hooks
		- xps hooks
		"""
		if not cfg.get("fallback_convert", self.cfg.get("fallback_convert", False)):
			return ""

		path = cfg.get("path") or self.state.get("path")
		if not path:
			return ""

		# Try LibreOffice/soffice -> txt
		soffice = shutil.which("soffice") or shutil.which("libreoffice")
		if soffice:
			created_tmp = False
			tmp_dir = cfg.get("tmp_dir") or self.cfg.get("tmp_dir")
			if not tmp_dir:
				tmp_dir = tempfile.mkdtemp(prefix="universaldoc_")
				created_tmp = True

			try:
				cmd = [
					soffice,
					"--headless",
					"--convert-to", "txt:Text",
					"--outdir", tmp_dir,
					path,
				]
				proc = subprocess.run(cmd, capture_output=True, text=True)
				if proc.returncode == 0:
					out_path = os.path.join(tmp_dir, os.path.splitext(os.path.basename(path))[0] + ".txt")
					if os.path.exists(out_path):
						with open(out_path, "r", encoding="utf-8", errors="replace") as f:
							return f.read()
			except Exception:
				pass
			finally:
				if created_tmp:
					shutil.rmtree(tmp_dir, ignore_errors=True)

		return ""

	def _infer_kind(self, cfg: Dict[str, Any]) -> str:
		kind = cfg.get("kind", "auto")
		path = cfg.get("path")

		if kind and kind != "auto":
			return str(kind).lower()

		if path:
			ext = os.path.splitext(path)[1].lower().strip(".")
			if ext:
				return ext

		return "text"

	def _effective_kind(self, cfg: Dict[str, Any]) -> str:
		return self._infer_kind({
			"path": cfg.get("path") or self.state.get("path"),
			"kind": cfg.get("kind", self.state.get("kind", "auto")),
		})

	def _clean_header(self, value: Any) -> str:
		s = "" if value is None else str(value).strip()
		if not s:
			return ""
		s = re.sub(r"[^\w\s]", "", s)
		parts = s.split()
		return "".join(p[:1].upper() + p[1:] for p in parts)

	def _merge_cfg(self, cfg: Optional[Dict[str, Any]], extra_defaults: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
		merged = dict(self.cfg)
		if extra_defaults:
			for k, v in extra_defaults.items():
				if k not in merged:
					merged[k] = v
		if cfg:
			merged.update(cfg)
		return merged

	def _defaults(self, cfg: Optional[Dict[str, Any]], defaults: Dict[str, Any]) -> Dict[str, Any]:
		out = dict(defaults)
		if cfg:
			out.update(cfg)
		return out


# =============================================================================
# Optional pretty printer module in same file
# =============================================================================

class UniversalDocPrinter:
	"""
	Pretty terminal renderer.

	Uses rich if available.
	Falls back to plain print if not.

	This is intentionally separate from the parsing class so you can swap
	renderers later without touching extraction logic.
	"""

	def __init__(self, cfg: Optional[Dict[str, Any]] = None):
		self.cfg = {
			"force_plain": False,
		}
		if cfg:
			self.cfg.update(cfg)

	def print_markdown(self, cfg: Optional[Dict[str, Any]] = None) -> None:
		cfg = cfg or {}
		text = cfg.get("text", "")

		if self.cfg.get("force_plain"):
			print(text)
			return

		try:
			from rich.console import Console
			from rich.markdown import Markdown
			console = Console()
			console.print(Markdown(text))
		except Exception:
			print(text)

	def print_blocks(self, cfg: Optional[Dict[str, Any]] = None) -> None:
		cfg = cfg or {}
		blocks = cfg.get("blocks", [])

		# Easiest path: render blocks to markdown-ish plain markdown
		tmp = UniversalDocLib({})
		md = tmp._blocks_to_markdown({"blocks": blocks})
		self.print_markdown({"text": md})

	def print_tables(self, cfg: Optional[Dict[str, Any]] = None) -> None:
		cfg = cfg or {}
		tables = cfg.get("tables", [])

		if self.cfg.get("force_plain"):
			for table in tables:
				print(" > ".join(table.get("path", [])))
				print(tmp._matrix_to_markdown({"matrix": [table.get("headers", [])] + table.get("rows", [])}))
				print()
			return

		try:
			from rich.console import Console
			from rich.table import Table
			from rich import box

			console = Console()
			for table_obj in tables:
				title = " > ".join(table_obj.get("path", [])) or "Table"
				headers = table_obj.get("headers", [])
				rows = table_obj.get("rows", [])

				table = Table(title=title, box=box.ROUNDED)
				for h in headers:
					table.add_column(str(h))
				for row in rows:
					table.add_row(*[str(x) for x in row])

				console.print(table)
				console.print()
		except Exception:
			tmp = UniversalDocLib({})
			for table_obj in tables:
				print(" > ".join(table_obj.get("path", [])))
				print(tmp._matrix_to_markdown({"matrix": [table_obj.get("headers", [])] + table_obj.get("rows", [])}))
				print()


# =============================================================================
# Example usage
# =============================================================================

if __name__ == "__main__":
	"""
	Basic quick test:

		python universal_doc_lib.py /path/to/file.docx

	Example import usage:

		from universal_doc_lib import UniversalDocLib, UniversalDocPrinter

		doc = UniversalDocLib({"path": "invoice.pdf"})
		md = doc.to_markdown({})
		blocks = doc.to_blocks({})
		tables = doc.to_tables({})

		printer = UniversalDocPrinter({})
		printer.print_markdown({"text": md})
		printer.print_tables({"tables": tables})
	"""
	import sys

	if len(sys.argv) > 1:
		path = sys.argv[1]
		doc = UniversalDocLib({"path": path})
		print(doc.to_markdown({}))
	else:
		print("Usage: python universal_doc_lib.py /path/to/file")